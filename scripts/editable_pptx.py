#!/usr/bin/env python3
"""
editable_pptx.py - 将AI生成的幻灯片PNG转为可编辑PPTX

直接调用 px-image2pptx 的核心组件：
- PaddleOCR PP-OCRv5 文字检测
- 经典CV文字掩码（双阈值 + OCR裁剪 + Canny边缘）
- LAMA 背景修复
- 精准文本框组装（自动字体缩放 + 颜色检测）

依赖：
  pip install px-image2pptx[ocr,inpaint]
"""

import argparse
import glob
import json
import os
import re
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple


# =============================================================================
# 核心：调用 px-image2pptx
# =============================================================================

def _convert_single_slide(
    image_path: str,
    output_pptx_path: str,
    skip_inpaint: bool = False,
    lang: str = "auto",
    work_dir: Optional[str] = None,
) -> Dict[str, Any]:
    """用 px-image2pptx 处理单张幻灯片"""
    from px_image2pptx.pipeline import image_to_pptx

    report = image_to_pptx(
        image_path=image_path,
        output_path=output_pptx_path,
        lang=lang,
        skip_inpaint=skip_inpaint,
        max_inpaint_size=1024,
        work_dir=work_dir,
    )
    return report


# =============================================================================
# 合并多页单页PPTX → 多页PPTX
# =============================================================================

def _merge_single_page_pptx(
    single_pptx_paths: List[str],
    merged_output_path: str,
    images_dir: str,
    slide_w: float = 13.333,
    slide_h: float = 7.5,
) -> str:
    """将多个单页PPTX合并为一个完整PPTX

    用 px-image2pptx 的 assemble 逐页重建，避免 python-pptx 的跨文档关系丢失。
    """
    import cv2
    from pptx import Presentation
    from pptx.util import Inches

    merged = Presentation()
    merged.slide_width = Inches(slide_w)
    merged.slide_height = Inches(slide_h)

    for src_path in single_pptx_paths:
        work_dir = src_path.replace("_editable.pptx", "_work")
        bg_path = os.path.join(work_dir, "background.png")
        ocr_json_path = os.path.join(work_dir, "text_regions.json")
        # 原图路径从文件名推断
        fname = os.path.basename(src_path).replace("_editable.pptx", ".png")
        img_candidates = [
            os.path.join(images_dir, fname),
            os.path.join(os.path.dirname(os.path.dirname(src_path)), "outputs", "images", fname),
            os.path.join(os.path.dirname(src_path), fname),
        ]
        img_path = None
        for c in img_candidates:
            if os.path.exists(c):
                img_path = c
                break

        if not img_path or not os.path.exists(ocr_json_path):
            # 回退：直接用单页pptx的shapes
            print(f"  跳过合并（缺少中间文件）: {src_path}")
            continue

        from px_image2pptx.ocr import load_ocr_json
        from px_image2pptx.assemble import (
            SlideMapper, group_text_lines, group_bbox, group_to_text,
            autoscale_font, detect_text_color, detect_bg_color,
        )
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from PIL import Image
        import numpy as np

        ocr_regions = load_ocr_json(ocr_json_path)
        img = Image.open(img_path)
        img_w, img_h = img.size
        mapper = SlideMapper(img_w, img_h)

        # 调整merged的尺寸为第一页的mapper
        if len(merged.slides) == 0:
            merged.slide_width = Inches(mapper.slide_w)
            merged.slide_height = Inches(mapper.slide_h)

        slide = merged.slides.add_slide(merged.slide_layouts[6])

        # 背景
        actual_bg = bg_path if os.path.exists(bg_path) else img_path
        slide.shapes.add_picture(
            actual_bg, Emu(0), Emu(0),
            Inches(mapper.slide_w), Inches(mapper.slide_h),
        )

        # 加载tight_mask（用于颜色检测）
        tight_mask_path = os.path.join(work_dir, "tight_mask.png")
        img_rgb = None
        tight_mask = None
        if os.path.exists(tight_mask_path):
            tight_mask = cv2.imread(tight_mask_path, cv2.IMREAD_GRAYSCALE)
            img_rgb = np.array(Image.open(img_path).convert("RGB"))

        # 文本框
        text_groups = group_text_lines(ocr_regions)
        for group in text_groups:
            x1, y1, x2, y2 = group_bbox(group)
            text = group_to_text(group)
            if not text.strip():
                continue

            left, top, width, height = mapper.bbox_to_emu(x1, y1, x2, y2)
            pad = mapper.to_emu(2)
            left = max(0, left - pad)
            top = max(0, top - pad)
            width += pad * 2
            height += pad * 2

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)

            bbox_w = x2 - x1
            region_heights = [r["bbox"]["y2"] - r["bbox"]["y1"] for r in group]
            line_h = sorted(region_heights)[len(region_heights) // 2]
            font_size = autoscale_font(text, bbox_w, line_h, mapper.ppi)

            if img_rgb is not None and tight_mask is not None:
                r, g, b = detect_text_color(img_rgb, tight_mask, x1, y1, x2, y2)
            else:
                r, g, b = 0x33, 0x33, 0x33
            color = RGBColor(r, g, b)

            lines = text.split("\n")
            p = tf.paragraphs[0]
            p.text = lines[0]
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            for line in lines[1:]:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(font_size)
                p.font.color.rgb = color

    merged.save(merged_output_path)
    return merged_output_path


# =============================================================================
# 完整流水线
# =============================================================================

def run_pipeline(
    images_dir: str,
    output_dir: str,
    slide_numbers: Optional[List[int]] = None,
    skip_inpaint: bool = False,
    title: str = "Untitled",
    lang: str = "auto",
) -> str:
    """完整流水线：遍历所有slide PNG → 可编辑PPTX（调用 px-image2pptx）"""
    # 收集图片
    pattern = os.path.join(images_dir, "slide-*.png")
    image_files = sorted(glob.glob(pattern))
    if not image_files:
        raise RuntimeError(f"未找到slide图片: {pattern}")

    # 过滤
    if slide_numbers:
        image_files = [
            f for f in image_files
            if int(os.path.basename(f).replace("slide-", "").replace(".png", "")) in slide_numbers
        ]

    print("=" * 60)
    print("可编辑PPTX转换流水线（px-image2pptx）")
    print("=" * 60)
    print(f"输入目录: {images_dir}")
    print(f"图片数量: {len(image_files)}")
    print(f"跳过修复: {skip_inpaint}")
    print("=" * 60)

    os.makedirs(output_dir, exist_ok=True)

    single_pptx_paths = []
    for img_path in image_files:
        fname = os.path.basename(img_path).replace(".png", "")
        slide_num = int(fname.replace("slide-", ""))
        out_pptx = os.path.join(output_dir, f"{fname}_editable.pptx")
        work_dir = os.path.join(output_dir, f"{fname}_work")

        print(f"\n  [slide {slide_num}] 处理 {os.path.basename(img_path)}...")

        report = _convert_single_slide(
            image_path=img_path,
            output_pptx_path=out_pptx,
            skip_inpaint=skip_inpaint,
            lang=lang,
            work_dir=work_dir,
        )

        print(f"  [slide {slide_num}] OCR区域: {report.get('ocr_regions', 0)}")
        print(f"  [slide {slide_num}] 文本框: {report.get('text_boxes', 0)}")
        print(f"  [slide {slide_num}] 背景: {report.get('background', 'unknown')}")
        timings = report.get("timings", {})
        print(f"  [slide {slide_num}] 耗时: OCR={timings.get('ocr', '?')}s, "
              f"textmask={timings.get('textmask', '?')}s, "
              f"inpaint={timings.get('inpaint', '?')}s, "
              f"assemble={timings.get('assemble', '?')}s")

        single_pptx_paths.append(out_pptx)

    # 合并为多页PPTX
    import re as _re
    safe_title = _re.sub(r"[^\w\u4e00-\u9fff\-]+", "_", title)[:60] or "deck"
    merged_path = os.path.join(output_dir, f"{safe_title}_editable.pptx")

    if len(single_pptx_paths) > 1:
        print(f"\n  合并 {len(single_pptx_paths)} 页...")
        _merge_single_page_pptx(single_pptx_paths, merged_path, images_dir)
        print(f"  📑 合并完成: {merged_path}")
    elif single_pptx_paths:
        merged_path = single_pptx_paths[0]

    print()
    print("=" * 60)
    print("转换完成！")
    print("=" * 60)
    print(f"输出目录: {output_dir}")
    print(f"可编辑PPTX: {merged_path}")

    return merged_path


# =============================================================================
# CLI
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="将AI生成的幻灯片PNG转为可编辑PPTX（基于 px-image2pptx）"
    )
    parser.add_argument("--images-dir", required=True, help="slide-XX.png 所在目录")
    parser.add_argument("--output", help="输出目录（默认：images-dir 的父目录）")
    parser.add_argument("--slides", help="只处理指定页，如 '1,3,5'")
    parser.add_argument("--skip-inpaint", action="store_true",
                        help="跳过LAMA背景修复（更快，但文本框下仍有原文字）")
    parser.add_argument("--title", default="Untitled", help="PPTX标题")
    parser.add_argument("--lang", default="auto",
                        help="OCR语言（auto=自动检测, ch=中英混合, en=纯英文）")

    args = parser.parse_args()

    slide_numbers = None
    if args.slides:
        slide_numbers = [int(x.strip()) for x in args.slides.split(",")]

    output_dir = args.output or os.path.dirname(args.images_dir)

    run_pipeline(
        images_dir=args.images_dir,
        output_dir=output_dir,
        slide_numbers=slide_numbers,
        skip_inpaint=args.skip_inpaint,
        title=args.title,
        lang=args.lang,
    )


if __name__ == "__main__":
    main()
